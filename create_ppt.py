from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import os

# === COLORS ===
BLUE_DARK   = RGBColor(0x1A, 0x3C, 0x6E)   # BPS dark blue
BLUE_MID    = RGBColor(0x00, 0x6E, 0xAF)   # BPS medium blue
BLUE_LIGHT  = RGBColor(0xDA, 0xE8, 0xFC)   # light blue bg
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT  = RGBColor(0xF5, 0xF5, 0xF5)
GRAY_TEXT   = RGBColor(0x44, 0x44, 0x44)
GREEN       = RGBColor(0x2E, 0x86, 0x48)
ORANGE      = RGBColor(0xE6, 0x7E, 0x22)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, W, Inches(1.3), fill=BLUE_DARK)
    add_text(slide, title, Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.7),
             size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(0.82), Inches(12.5), Inches(0.4),
                 size=14, color=RGBColor(0xAA, 0xCC, 0xFF), align=PP_ALIGN.LEFT)

def img_placeholder(slide, label, x, y, w, h):
    add_rect(slide, x, y, w, h, fill=BLUE_LIGHT, line=BLUE_MID, line_w=Pt(1.5))
    add_text(slide, f"[ INSERT IMAGE: {label} ]", x, y + h/2 - Inches(0.3),
             w, Inches(0.6), size=14, bold=True,
             color=BLUE_MID, align=PP_ALIGN.CENTER)

def bullet_box(slide, items, x, y, w, h, size=16, head=None, head_size=18):
    if head:
        add_text(slide, head, x, y, w, Inches(0.45),
                 size=head_size, bold=True, color=BLUE_DARK)
        y += Inches(0.45)
        h -= Inches(0.45)
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = item
        run.font.size  = Pt(size)
        run.font.color.rgb = GRAY_TEXT

# ─────────────────────────────────────────────────────────────
# SLIDE 1 — COVER
# ─────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=BLUE_DARK)
add_rect(sl, 0, Inches(5.8), W, Inches(1.7), fill=BLUE_MID)

add_text(sl, "SISTEM CHATBOT RAG BPS PEKANBARU",
         Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.9),
         size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Berbasis WhatsApp & Kecerdasan Buatan",
         Inches(0.8), Inches(2.1), Inches(11.7), Inches(0.7),
         size=24, color=RGBColor(0xAA, 0xCC, 0xFF), align=PP_ALIGN.CENTER)
add_rect(sl, Inches(5.4), Inches(2.9), Inches(2.5), Inches(0.05), fill=BLUE_LIGHT)

add_text(sl, "Menggunakan: n8n  ·  Gemini AI  ·  Supabase pgvector  ·  WAHA  ·  React",
         Inches(0.8), Inches(3.1), Inches(11.7), Inches(0.5),
         size=16, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)

add_text(sl, "Badan Pusat Statistik Kota Pekanbaru  ·  2025",
         Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.5),
         size=14, color=WHITE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────
# SLIDE 2 — AGENDA
# ─────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=GRAY_LIGHT)
header_bar(sl, "Agenda Presentasi")

items = [
    ("01", "Latar Belakang & Tujuan"),
    ("02", "Arsitektur Sistem"),
    ("03", "Alur Kerja Chatbot (n8n Workflow)"),
    ("04", "Use Case Diagram"),
    ("05", "ERD Database (Supabase)"),
    ("06", "Admin Dashboard (localhost:5000)"),
    ("07", "Tech Stack & Layanan"),
    ("08", "Cara Penggunaan & Demo"),
    ("09", "Kesimpulan"),
]
cols = [items[:5], items[5:]]
for ci, col in enumerate(cols):
    cx = Inches(0.5 + ci * 6.5)
    for ri, (num, txt) in enumerate(col):
        cy = Inches(1.5 + ri * 1.1)
        add_rect(sl, cx, cy, Inches(0.6), Inches(0.7), fill=BLUE_MID)
        add_text(sl, num, cx, cy + Inches(0.1), Inches(0.6), Inches(0.6),
                 size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, txt, cx + Inches(0.7), cy + Inches(0.1),
                 Inches(5.5), Inches(0.6), size=16, color=GRAY_TEXT)

# ─────────────────────────────────────────────────────────────
# SLIDE 3 — LATAR BELAKANG
# ─────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=WHITE)
header_bar(sl, "Latar Belakang", "Mengapa Chatbot RAG dibutuhkan BPS Pekanbaru?")

problems = [
    "📞  Staf BPS sering menerima pertanyaan statistik berulang dari masyarakat",
    "📄  Publikasi BPS (PDF) sulit dicari secara manual oleh pengguna awam",
    "⏰  Pertanyaan datang di luar jam kerja — tidak ada yang menjawab",
    "📊  Data BPS tersebar di banyak publikasi, sulit ditemukan secara cepat",
]
solutions = [
    "✅  Chatbot WhatsApp otomatis — menjawab 24/7 tanpa intervensi manual",
    "✅  RAG (Retrieval-Augmented Generation) — jawaban dari dokumen resmi BPS",
    "✅  Gemini AI — memahami pertanyaan dalam bahasa natural (Indonesia)",
    "✅  Real-time BPS API — data terkini langsung dari webapi.bps.go.id",
]
add_rect(sl, Inches(0.3), Inches(1.4), Inches(6.1), Inches(5.7), fill=RGBColor(0xFF,0xF0,0xF0))
add_text(sl, "MASALAH", Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.4),
         size=15, bold=True, color=RGBColor(0xC0,0x39,0x2B))
bullet_box(sl, problems, Inches(0.4), Inches(1.95), Inches(5.9), Inches(4.8), size=14)

add_rect(sl, Inches(6.7), Inches(1.4), Inches(6.3), Inches(5.7), fill=RGBColor(0xF0,0xFF,0xF0))
add_text(sl, "SOLUSI", Inches(6.9), Inches(1.5), Inches(6.0), Inches(0.4),
         size=15, bold=True, color=GREEN)
bullet_box(sl, solutions, Inches(6.8), Inches(1.95), Inches(6.1), Inches(4.8), size=14)

# ─────────────────────────────────────────────────────────────
# SLIDE 4 — ARSITEKTUR SISTEM
# ─────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=WHITE)
header_bar(sl, "Arsitektur Sistem", "Gambaran keseluruhan komponen dan alur data")
img_placeholder(sl, "Arsitektur Sistem — arsitektur.png",
                Inches(0.4), Inches(1.45), Inches(12.5), Inches(5.8))

# ─────────────────────────────────────────────────────────────
# SLIDE 5 — ALUR CHATBOT
# ─────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=WHITE)
header_bar(sl, "Alur Kerja Chatbot (n8n Workflow)", "Dari pesan WhatsApp hingga jawaban AI")
img_placeholder(sl, "Flowchart Chatbot n8n — flowchart_n8n.png",
                Inches(0.4), Inches(1.45), Inches(8.0), Inches(5.8))
# Side notes
notes = [
    "1. User kirim pesan WhatsApp",
    "2. WAHA terima & teruskan ke n8n",
    "3. AI Agent (Gemini) proses",
    "4. Cari di Knowledge Base (RAG)",
    "5. Jika perlu → query BPS API",
    "6. Jawaban dikirim balik ke WA",
]
bullet_box(sl, notes, Inches(8.6), Inches(1.6), Inches(4.4), Inches(5.5), size=15)

# ─────────────────────────────────────────────────────────────
# SLIDE 6 — USE CASE
# ─────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=WHITE)
header_bar(sl, "Use Case Diagram", "Aktor dan fungsionalitas sistem")
img_placeholder(sl, "Use Case Diagram — usecase.png",
                Inches(0.4), Inches(1.45), Inches(12.5), Inches(5.8))

# ─────────────────────────────────────────────────────────────
# SLIDE 7 — ERD
# ─────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=WHITE)
header_bar(sl, "Entity Relationship Diagram (ERD)", "Skema database Supabase PostgreSQL + pgvector")
img_placeholder(sl, "ERD Database — erd_dbdiagram.png",
                Inches(0.4), Inches(1.45), Inches(8.5), Inches(5.8))

tbl_info = [
    ("📚 documents", "Menyimpan chunk teks PDF\n+ embedding vector(768)"),
    ("📄 ingested_files", "Metadata file PDF\nyang sudah diproses"),
    ("💬 n8n_chat_histories", "Riwayat percakapan\nWhatsApp per sesi"),
]
for i, (name, desc) in enumerate(tbl_info):
    cy = Inches(1.55 + i * 1.85)
    add_rect(sl, Inches(9.1), cy, Inches(3.9), Inches(1.7), fill=BLUE_LIGHT, line=BLUE_MID, line_w=Pt(1))
    add_text(sl, name, Inches(9.2), cy + Inches(0.1), Inches(3.7), Inches(0.45),
             size=13, bold=True, color=BLUE_DARK)
    add_text(sl, desc, Inches(9.2), cy + Inches(0.55), Inches(3.7), Inches(1.0),
             size=12, color=GRAY_TEXT)

# ─────────────────────────────────────────────────────────────
# SLIDE 8 — ADMIN DASHBOARD FLOW
# ─────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=WHITE)
header_bar(sl, "Alur Admin Dashboard (localhost:5000)", "3 menu utama: Dashboard · PDF Processor · Credentials")
img_placeholder(sl, "Flowchart Admin Dashboard — flowchart_admin.png",
                Inches(0.4), Inches(1.45), Inches(12.5), Inches(5.8))

# ─────────────────────────────────────────────────────────────
# SLIDE 9 — TECH STACK
# ─────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=WHITE)
header_bar(sl, "Tech Stack & Layanan", "Komponen teknologi yang digunakan")

stack = [
    ("🤖  AI & Model",       ["Google Gemini 2.0 Flash (Chat)",
                               "Gemini Embedding 004 (768-dim)",
                               "RAG (Retrieval-Augmented Generation)"]),
    ("⚙️  Workflow & API",   ["n8n (Workflow Automation, port 5678)",
                               "WAHA (WhatsApp Gateway, port 3001)",
                               "BPS Web API (webapi.bps.go.id)"]),
    ("🗄️  Database",          ["Supabase PostgreSQL + pgvector",
                               "HNSW Index (cosine similarity)",
                               "match_documents RPC Function"]),
    ("🖥️  Frontend & Backend",["React 18 + Vite + Tailwind CSS",
                               "FastAPI (Python, port 8503)",
                               "Admin Dashboard (port 5000)"]),
    ("🐳  Infrastructure",   ["Docker (WAHA, Metabase)",
                               "Metabase Analytics (port 3002)",
                               "Windows (start.bat launcher)"]),
]
cols_per_row = 3
for i, (title, items) in enumerate(stack):
    row = i // cols_per_row
    col = i % cols_per_row
    cx = Inches(0.3 + col * 4.35)
    cy = Inches(1.45 + row * 2.8)
    add_rect(sl, cx, cy, Inches(4.1), Inches(2.6), fill=BLUE_LIGHT, line=BLUE_MID, line_w=Pt(1))
    add_text(sl, title, cx + Inches(0.1), cy + Inches(0.1), Inches(3.9), Inches(0.45),
             size=14, bold=True, color=BLUE_DARK)
    bullet_box(sl, ["• " + x for x in items],
               cx + Inches(0.1), cy + Inches(0.55), Inches(3.9), Inches(2.0), size=13)

# ─────────────────────────────────────────────────────────────
# SLIDE 10 — CARA PENGGUNAAN
# ─────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=WHITE)
header_bar(sl, "Cara Penggunaan Sistem", "Operasional harian dan fitur utama")

steps_user = [
    "1. Kirim pesan ke nomor WhatsApp bot",
    "2. Tanyakan data statistik BPS Pekanbaru",
    "3. Bot menjawab otomatis dalam hitungan detik",
    "4. Jawaban berdasarkan publikasi resmi BPS",
]
steps_admin = [
    "1. Buka localhost:5000 (Admin Dashboard)",
    "2. Upload PDF publikasi BPS baru",
    "3. Sistem otomatis ekstrak, chunk, dan embed",
    "4. Kelola API credentials di menu Credentials",
]
add_rect(sl, Inches(0.3), Inches(1.45), Inches(6.1), Inches(5.7), fill=RGBColor(0xE8,0xF4,0xFF))
add_text(sl, "👤  Untuk Pengguna WhatsApp",
         Inches(0.5), Inches(1.55), Inches(5.7), Inches(0.45),
         size=16, bold=True, color=BLUE_DARK)
bullet_box(sl, steps_user, Inches(0.5), Inches(2.05), Inches(5.7), Inches(4.8), size=15)

add_rect(sl, Inches(6.9), Inches(1.45), Inches(6.1), Inches(5.7), fill=RGBColor(0xF0,0xF8,0xE8))
add_text(sl, "🔧  Untuk Admin BPS",
         Inches(7.1), Inches(1.55), Inches(5.7), Inches(0.45),
         size=16, bold=True, color=GREEN)
bullet_box(sl, steps_admin, Inches(7.1), Inches(2.05), Inches(5.7), Inches(4.8), size=15)

# ─────────────────────────────────────────────────────────────
# SLIDE 11 — LAYANAN & PORT
# ─────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=WHITE)
header_bar(sl, "Layanan & Port", "Ringkasan semua service yang berjalan")

services = [
    ("Admin Dashboard", "React + Vite",    "5000",  "start.bat",             BLUE_MID),
    ("PDF Processor",   "FastAPI Python",  "8503",  "start.bat",             BLUE_MID),
    ("n8n Workflow",    "Node.js",         "5678",  "n8n start (manual)",    ORANGE),
    ("WAHA Gateway",    "Docker",          "3001",  "Docker Desktop",        GREEN),
    ("Metabase",        "Docker",          "3002",  "docker compose up -d",  GRAY_TEXT),
    ("Supabase",        "Cloud (PG+pgvec)","—",     "Sudah berjalan cloud",  RGBColor(0x6A,0x0D,0xAD)),
]
# header row
cols_w = [Inches(2.6), Inches(2.6), Inches(1.4), Inches(3.0), Inches(3.0)]
cx_list = [Inches(0.3), Inches(2.9), Inches(5.5), Inches(6.9), Inches(9.9)]
headers = ["Layanan", "Teknologi", "Port", "Cara Jalankan", "Keterangan"]
for ci, (hdr, cx, cw) in enumerate(zip(headers, cx_list, cols_w)):
    add_rect(sl, cx, Inches(1.45), cw, Inches(0.5), fill=BLUE_DARK)
    add_text(sl, hdr, cx + Inches(0.05), Inches(1.5), cw, Inches(0.4),
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

keterangan = ["Antarmuka upload & kelola PDF",
              "Backend pemrosesan PDF",
              "Orkestrasi AI chatbot",
              "WhatsApp gateway",
              "Analytics (opsional)",
              "Database vektor cloud"]
for ri, (name, tech, port, cara, color) in enumerate(services):
    cy = Inches(1.95 + ri * 0.82)
    bg = BLUE_LIGHT if ri % 2 == 0 else WHITE
    add_rect(sl, Inches(0.3), cy, Inches(12.7), Inches(0.82), fill=bg)
    vals = [name, tech, port, cara, keterangan[ri]]
    for ci, (val, cx, cw) in enumerate(zip(vals, cx_list, cols_w)):
        clr = color if ci == 0 else GRAY_TEXT
        add_text(sl, val, cx + Inches(0.05), cy + Inches(0.1), cw, Inches(0.6),
                 size=12, bold=(ci == 0), color=clr)

# ─────────────────────────────────────────────────────────────
# SLIDE 12 — KESIMPULAN
# ─────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=BLUE_DARK)
add_rect(sl, 0, Inches(5.5), W, Inches(2.0), fill=BLUE_MID)

add_text(sl, "Kesimpulan", Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.7),
         size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

points = [
    "✅  Chatbot WhatsApp otomatis berbasis RAG berhasil dibangun untuk BPS Pekanbaru",
    "✅  Sistem dapat menjawab pertanyaan statistik 24/7 dari dokumen resmi BPS",
    "✅  Admin dashboard memudahkan pengelolaan PDF tanpa keahlian teknis",
    "✅  Integrasi BPS Web API memastikan data real-time tersedia",
    "✅  Teknologi modern: Gemini AI + pgvector + n8n + WAHA",
]
for i, pt in enumerate(points):
    add_text(sl, pt, Inches(1.5), Inches(1.3 + i * 0.78), Inches(10.3), Inches(0.7),
             size=16, color=RGBColor(0xCC, 0xDD, 0xFF))

add_text(sl, "Terima Kasih",
         Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.8),
         size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Badan Pusat Statistik Kota Pekanbaru  ·  2025",
         Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5),
         size=14, color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────
out = r"C:\MachineLearning Project\BPS-n8n-RAG_ChatBot\BPS_RAG_Chatbot_Presentasi.pptx"
prs.save(out)
print(f"PPT saved: {out}")
print(f"Total slides: {len(prs.slides)}")
